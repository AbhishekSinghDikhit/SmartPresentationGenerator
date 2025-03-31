from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

TEMPLATE_DIR = "templates"
os.makedirs(TEMPLATE_DIR, exist_ok=True)

def set_slide_background(slide, prs, bg_color=None, bg_image=None):
    """Sets the background of a slide using either a solid color or an image."""
    if bg_image and os.path.exists(bg_image):
        # Add a full-size background image
        slide.shapes.add_picture(bg_image, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    elif isinstance(bg_color, RGBColor):  # Ensure bg_color is a valid RGBColor
        # Apply solid color background
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color
    else:
        print("⚠️ Invalid bg_color! It must be an RGBColor object.")

def create_template(template_name: str, title_color, bg_color, text_color, accent_color, bg_image=None, title_text="Presentation Title", author_name="Author Name", bullets=None, img_placeholder_path="slide_images/placeholder.jpg"):
    """Creates a PowerPoint template with multiple layouts and saves it in the templates folder."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ---- TITLE SLIDE ----
    title_slide_layout = prs.slide_layouts[5]  # Use Title Slide Layout
    slide = prs.slides.add_slide(title_slide_layout)
    set_slide_background(slide, prs, bg_color, bg_image)

    # Title
    if slide.shapes.title:
        slide.shapes.title.text = title_text
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = title_color

    # Author Name
    author_shape = slide.shapes.add_textbox(Inches(4), Inches(5), Inches(5), Inches(1))
    author_frame = author_shape.text_frame
    author_frame.text = f"By {author_name}"
    author_frame.paragraphs[0].font.size = Pt(28)
    author_frame.paragraphs[0].font.color.rgb = title_color

    # ---- CONTENT SLIDE ----
    content_slide_layout = prs.slide_layouts[1]  # Title & Content Layout
    slide = prs.slides.add_slide(content_slide_layout)
    set_slide_background(slide, prs, bg_color, bg_image)

    # Set Title
    if slide.shapes.title:
        slide.shapes.title.text = "Content Slide Title"
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = text_color

    # Bullet Points
    if bullets:
        content_box = slide.shapes.placeholders[1].text_frame
        for bullet in bullets:
            p = content_box.add_paragraph()
            p.text = bullet
            p.font.size = Pt(24)
            p.font.color.rgb = text_color

    # ---- IMAGE + TEXT SLIDE ----
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank Layout
    set_slide_background(slide, prs, bg_color, bg_image)

    # Image Placeholder
    img_placeholder = slide.shapes.add_picture(img_placeholder_path, Inches(1), Inches(1.5), width=Inches(5), height=Inches(4))

    # Text Box
    text_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5), Inches(4))
    text_frame = text_box.text_frame
    text_frame.text = "Image + Text Layout"
    text_frame.paragraphs[0].font.size = Pt(32)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = accent_color

    # ---- THANK YOU SLIDE ----
    slide = prs.slides.add_slide(title_slide_layout)
    set_slide_background(slide, prs, bg_color, bg_image)

    if slide.shapes.title:
        slide.shapes.title.text = "Thank You!"
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(48)
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = title_color
        slide.shapes.title.text_frame.paragraphs[0].alignment = 1  # Center

    # Save Template
    file_path = os.path.join(TEMPLATE_DIR, f"{template_name}.pptx")
    prs.save(file_path)
    print(f"✅ Template saved: {file_path}")

# Example usage
create_template(
    template_name="modern_minimalist",
    title_color=RGBColor(255, 255, 255),
    bg_color=RGBColor(0, 51, 102),
    text_color=RGBColor(0, 0, 0),
    accent_color=RGBColor(255, 87, 34),
    title_text="Dynamic Presentation Title",
    author_name="John Doe",
    bullets=["Point 1", "Point 2", "Point 3"]
)

create_template(
    template_name="corporate_professional",
    title_color=RGBColor(0, 0, 0),
    bg_color=RGBColor(220, 220, 220),
    text_color=RGBColor(50, 50, 50),
    accent_color=RGBColor(0, 102, 204),
    bg_image="slide_images/corporate_bg.jpg",
    title_text="Corporate Strategy",
    author_name="Jane Smith",
    bullets=["Business Goal 1", "Business Goal 2"]
)
