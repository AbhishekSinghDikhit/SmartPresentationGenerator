from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE
import os
from typing import List, Dict
from dotenv import load_dotenv
import requests
from PIL import Image
import re

load_dotenv()

TEMPLATE_DIR = "templates"
IMAGE_DIR = "slide_images"
os.makedirs(IMAGE_DIR, exist_ok=True)

def sanitize_filename(filename: str) -> str:
    """Removes invalid characters from filenames."""
    return re.sub(r'[<>:"/\\|?*]', '_', filename)

def get_template_path(template_name: str) -> str:
    """Returns the full path of the selected template."""
    template_path = os.path.join(TEMPLATE_DIR, f"{template_name}.pptx")
    if os.path.exists(template_path):
        return template_path
    else:
        raise FileNotFoundError(f"❌ Template '{template_name}' not found!")

def generate_slide_image(prompt: str, style: str) -> str:
    """
    Generates an image using the Cloudflare Worker API with the given style.
    """
    try:
        WORKER_URL = "https://image-generator.worldforscience.workers.dev/"
        response = requests.post(WORKER_URL, json={"prompt": prompt, "style": style})
        
        if response.status_code == 200:
            image_data = response.content
            image_path = os.path.join(IMAGE_DIR, f"{sanitize_filename(prompt)}_{style}.png")

            with open(image_path, "wb") as img_file:
                img_file.write(image_data)

            return image_path
        else:
            print(f"❌ Error: {response.status_code} - {response.text}")
            return None

    except Exception as e:
        print(f"❌ Error generating image with Cloudflare Worker: {e}")
        return None

def adjust_font_size(text_frame, max_lines=8, max_font_size=Pt(20), min_font_size=Pt(12)):
    """
    Dynamically adjusts font size to fit content within the text box.
    """
    while text_frame.text and len(text_frame.text.split("\n")) > max_lines and max_font_size > min_font_size:
        max_font_size -= Pt(2)
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = max_font_size

def add_chart_to_slide(slide, chart_data: Dict[str, List[float]], chart_type: str = "bar") -> None:
    """
    Adds a chart (bar, pie, etc.) to the slide based on the provided data.
    """
    chart_data_obj = CategoryChartData()
    chart_data_obj.categories = chart_data["categories"]
    chart_data_obj.add_series('Series 1', chart_data["values"])

    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(5), Inches(3)  # Position and size of the chart
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data_obj
    ).chart

    chart.has_legend = True
    chart.plots[0].has_data_labels = True

def add_shape_to_slide(slide, shape_type: str, left: Inches, top: Inches, width: Inches, height: Inches) -> None:
    """
    Adds a shape (rectangle, oval, etc.) to the slide.
    """
    if shape_type == "rectangle":
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
    elif shape_type == "circle":
        slide.shapes.add_shape(
            MSO_SHAPE.ELLIPSE, left, top, width, height
        )

def generate_pptx(request, ppt_content: List[Dict[str, List[str]]], image_style: str = "realistic", template_name: str = None) -> str:
    try:
        # Load selected template
        template_path = get_template_path(template_name or request.template)
        prs = Presentation(template_path)

        # Ensure template has at least 2 slides (Title + Thank You)
        if len(prs.slides) < 2:
            raise ValueError("Template must have at least a Title slide and a Thank You slide.")

        # Title Slide (First slide from template)
        title_slide = prs.slides[0]  
        if title_slide.shapes.title:
            title_slide.shapes.title.text = request.title.strip()

        for shape in title_slide.shapes:
            if shape.has_text_frame and "By" in shape.text:
                shape.text_frame.text = f"By {request.author.strip()}"
                break

        # Content Slides (Ensuring the number of slides does not exceed the template's layout count)
        available_slides = min(len(ppt_content), request.num_slides)
        left_side = True  # Toggle image placement

        for i, slide_data in enumerate(ppt_content[:available_slides]):
            slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]  # Use a content layout
            slide = prs.slides.add_slide(slide_layout)  

            # Set Slide Title
            if slide.shapes.title:
                slide.shapes.title.text = slide_data["title"].strip()
                slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
                slide.shapes.title.text_frame.paragraphs[0].font.bold = True
                slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

            # Generate Image
            image_path = generate_slide_image(slide_data["title"], image_style)
            img_width, img_height = Inches(5), Inches(4)

            if image_path:
                img_left = Inches(0.5) if left_side else Inches(7.5)
                text_left = Inches(7.5) if left_side else Inches(0.5)
                left_side = not left_side  # Toggle for next slide

                slide.shapes.add_picture(image_path, img_left, Inches(1.5), width=img_width, height=img_height)
                text_width = Inches(5.5)
            else:
                text_left = Inches(1)
                text_width = Inches(10)

            # Text Box for Content
            text_top, text_height = Inches(1.5), Inches(5)
            content_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            text_frame.margin_bottom = Inches(0.2)

            # Add Text Bullets
            font_size = Pt(20)
            for bullet in slide_data["content"]:
                p = text_frame.add_paragraph()
                p.text = bullet.strip()
                p.font.size = font_size
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.level = 0

            # Adjust font size if text overflows
            adjust_font_size(text_frame)

            # Add chart if specified in slide_data
            if "chart" in slide_data:
                add_chart_to_slide(slide, slide_data["chart"])

            # Add shapes if specified in slide_data
            if "shapes" in slide_data:
                for shape in slide_data["shapes"]:
                    add_shape_to_slide(slide, shape["type"], Inches(shape["left"]), Inches(shape["top"]), Inches(shape["width"]), Inches(shape["height"]))

        # Remove the last slide ("Thank You") and add it at the end
        thank_you_slide = prs.slides[-1]  # Assuming the last slide is "Thank You"
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])  # Remove "Thank You" slide

        # Add "Thank You" slide at the end
        thank_you_layout = prs.slide_layouts[5]  # Title slide layout for "Thank You"
        slide = prs.slides.add_slide(thank_you_layout)
        slide.shapes.title.text = "Thank You!"
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Save Presentation
        file_name = f"{sanitize_filename(request.title.strip())}.pptx"
        file_path = os.path.join(os.getcwd(), file_name)
        prs.save(file_path)
        return file_path

    except PermissionError:
        raise OSError(f"❌ Permission denied when saving file to {file_path}")
    except Exception as e:
        raise Exception(f"❌ Failed to generate presentation: {str(e)}")
